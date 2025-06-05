#!/usr/bin/env python
# -*- coding: utf-8 -*-

"""
增强版PPT生成脚本 - 图文并茂的专业演示文稿
基于论文内容，加入更多可视化元素和技术细节
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
import os

def add_gradient_background(slide, color1, color2):
    """添加渐变背景"""
    background = slide.background
    fill = background.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = color1
    fill.gradient_stops[1].color.rgb = color2

def create_title_slide(prs):
    """创建更加精美的标题页，加入图形元素"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 深蓝色渐变背景
    add_gradient_background(slide, RGBColor(25, 52, 65), RGBColor(62, 92, 118))
    
    shapes = slide.shapes
    
    # 添加装饰性几何图形
    # 左上角装饰
    for i in range(3):
        circle = shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(0.5 + i * 0.3), Inches(0.5 + i * 0.3),
            Inches(0.5 - i * 0.1), Inches(0.5 - i * 0.1)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = RGBColor(52, 152, 219)
        circle.fill.transparency = 0.3 + i * 0.2
        circle.line.fill.background()
    
    # 右下角装饰
    for i in range(3):
        rect = shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(8.5 - i * 0.3), Inches(4.5 - i * 0.3),
            Inches(0.8), Inches(0.5)
        )
        rect.fill.solid()
        rect.fill.fore_color.rgb = RGBColor(46, 204, 113)
        rect.fill.transparency = 0.3 + i * 0.2
        rect.line.fill.background()
    
    # 中央标题区域
    title_bg = shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1), Inches(1.5),
        Inches(8), Inches(2.5)
    )
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
    title_bg.fill.transparency = 0.1
    title_bg.line.color.rgb = RGBColor(52, 152, 219)
    title_bg.line.width = Pt(3)
    
    # 主标题
    title_box = shapes.add_textbox(Inches(1.2), Inches(1.7), Inches(7.6), Inches(2))
    title_frame = title_box.text_frame
    title_frame.clear()
    p = title_frame.paragraphs[0]
    p.text = u"基于Raft算法的"
    p.font.name = u"微软雅黑"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    p = title_frame.add_paragraph()
    p.text = u"分布式数据库构建"
    p.font.name = u"微软雅黑"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # 英文副标题
    p = title_frame.add_paragraph()
    p.text = u"Distributed Database Construction Based on Raft Algorithm"
    p.font.size = Pt(18)
    p.font.italic = True
    p.font.color.rgb = RGBColor(189, 195, 199)
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(12)
    
    # 学校和院系信息
    dept_box = shapes.add_textbox(Inches(3), Inches(0.2), Inches(4), Inches(0.4))
    dept_frame = dept_box.text_frame
    p = dept_frame.paragraphs[0]
    p.text = u"计算机科学与技术学院"
    p.font.name = u"微软雅黑"
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # 答辩信息
    info_box = shapes.add_textbox(Inches(2), Inches(4.2), Inches(6), Inches(1))
    info_frame = info_box.text_frame
    info_frame.clear()
    p = info_frame.paragraphs[0]
    p.text = u"答辩人：张皓然    学号：202131061326"
    p.font.name = u"微软雅黑"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    p = info_frame.add_paragraph()
    p.text = u"指导教师：蒋欣岑 讲师"
    p.font.name = u"微软雅黑"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(236, 240, 241)
    p.alignment = PP_ALIGN.CENTER

def create_outline_slide(prs):
    """创建动态目录页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shapes = slide.shapes
    
    # 浅色背景
    add_gradient_background(slide, RGBColor(245, 247, 250), RGBColor(255, 255, 255))
    
    # 标题
    title_shape = shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(2.5), Inches(0.3),
        Inches(5), Inches(0.8)
    )
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = RGBColor(41, 128, 185)
    title_shape.line.fill.background()
    
    title_box = shapes.add_textbox(Inches(2.5), Inches(0.4), Inches(5), Inches(0.6))
    text_frame = title_box.text_frame
    p = text_frame.paragraphs[0]
    p.text = u"演示大纲"
    p.font.name = u"微软雅黑"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # 目录项
    outline_items = [
        (u"01", u"研究背景与意义", u"大数据时代的挑战与机遇"),
        (u"02", u"理论基础", u"CAP定理、BASE理论、共识算法"),
        (u"03", u"Raft算法详解", u"核心机制与实现原理"),
        (u"04", u"系统架构设计", u"分层架构与模块设计"),
        (u"05", u"关键技术实现", u"日志复制、Leader选举、故障恢复"),
        (u"06", u"实验结果分析", u"功能测试与性能评估"),
        (u"07", u"总结与展望", u"研究成果与未来方向")
    ]
    
    for i, (num, title, subtitle) in enumerate(outline_items):
        # 创建卡片式布局
        card = shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(1.2 + (i % 2) * 4.3), Inches(1.5 + (i // 2) * 1),
            Inches(3.8), Inches(0.8)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(255, 255, 255)
        card.line.color.rgb = RGBColor(52, 152, 219)
        card.line.width = Pt(2)
        
        # 编号圆圈
        num_circle = shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(1.3 + (i % 2) * 4.3), Inches(1.55 + (i // 2) * 1),
            Inches(0.7), Inches(0.7)
        )
        num_circle.fill.solid()
        num_circle.fill.fore_color.rgb = RGBColor(52, 152, 219)
        num_circle.line.fill.background()
        
        # 编号文字
        num_text = shapes.add_textbox(
            Inches(1.3 + (i % 2) * 4.3), Inches(1.7 + (i // 2) * 1),
            Inches(0.7), Inches(0.4)
        )
        tf = num_text.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.name = u"Arial"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        # 标题和副标题
        text_box = shapes.add_textbox(
            Inches(2.1 + (i % 2) * 4.3), Inches(1.55 + (i // 2) * 1),
            Inches(2.8), Inches(0.7)
        )
        tf = text_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = u"微软雅黑"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RGBColor(44, 62, 80)
        
        p = tf.add_paragraph()
        p.text = subtitle
        p.font.name = u"微软雅黑"
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(127, 140, 141)

def create_background_slide(prs):
    """创建研究背景幻灯片，包含数据增长图表"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shapes = slide.shapes
    
    # 添加标题
    title_bg = shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(1)
    )
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = RGBColor(41, 128, 185)
    title_bg.line.fill.background()
    
    title_box = shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = u"研究背景 - 大数据时代的挑战"
    p.font.name = u"微软雅黑"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # 左侧文字内容
    left_box = shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4.5), Inches(3.5))
    text_frame = left_box.text_frame
    text_frame.clear()
    
    # 数据爆炸式增长
    p = text_frame.paragraphs[0]
    p.text = u"数据爆炸式增长"
    p.font.name = u"微软雅黑"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(41, 128, 185)
    
    points = [
        u"全球数据量以指数级增长",
        u"2025年预计达到175ZB",
        u"90%的数据在过去2年产生"
    ]
    
    for point in points:
        p = text_frame.add_paragraph()
        p.text = u"▸ " + point
        p.font.name = u"微软雅黑"
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(52, 73, 94)
        p.space_before = Pt(8)
        p.level = 1
    
    # 传统数据库的瓶颈
    p = text_frame.add_paragraph()
    p.text = u"\n传统数据库的瓶颈"
    p.font.name = u"微软雅黑"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(41, 128, 185)
    p.space_before = Pt(20)
    
    bottlenecks = [
        u"单机存储容量限制",
        u"并发处理能力不足",
        u"扩展成本高昂"
    ]
    
    for bottleneck in bottlenecks:
        p = text_frame.add_paragraph()
        p.text = u"▸ " + bottleneck
        p.font.name = u"微软雅黑"
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(52, 73, 94)
        p.space_before = Pt(8)
        p.level = 1
    
    # 右侧添加图表
    chart_data = ChartData()
    chart_data.categories = ['2015', '2017', '2019', '2021', '2023', '2025(预测)']
    chart_data.add_series('全球数据量(ZB)', (10, 25, 45, 79, 120, 175))
    
    x, y, cx, cy = Inches(5.5), Inches(1.5), Inches(4), Inches(3)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    ).chart
    
    # 设置图表样式
    chart.has_title = True
    chart.chart_title.text_frame.text = u"全球数据量增长趋势"
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(16)
    chart.chart_title.text_frame.paragraphs[0].font.bold = True

def create_cap_slide(prs):
    """创建CAP理论图解幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shapes = slide.shapes
    
    # 标题
    title_bg = shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(1)
    )
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = RGBColor(41, 128, 185)
    title_bg.line.fill.background()
    
    title_box = shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = u"CAP理论 - 分布式系统的基本约束"
    p.font.name = u"微软雅黑"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # 创建CAP三角形图
    # C顶点
    c_circle = shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(4.5), Inches(1.5),
        Inches(1), Inches(1)
    )
    c_circle.fill.solid()
    c_circle.fill.fore_color.rgb = RGBColor(231, 76, 60)
    c_circle.line.fill.background()
    
    c_text = shapes.add_textbox(Inches(4.5), Inches(1.7), Inches(1), Inches(0.6))
    tf = c_text.text_frame
    p = tf.paragraphs[0]
    p.text = "C"
    p.font.name = "Arial"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # A顶点
    a_circle = shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(2), Inches(3.5),
        Inches(1), Inches(1)
    )
    a_circle.fill.solid()
    a_circle.fill.fore_color.rgb = RGBColor(46, 204, 113)
    a_circle.line.fill.background()
    
    a_text = shapes.add_textbox(Inches(2), Inches(3.7), Inches(1), Inches(0.6))
    tf = a_text.text_frame
    p = tf.paragraphs[0]
    p.text = "A"
    p.font.name = "Arial"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # P顶点
    p_circle = shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(7), Inches(3.5),
        Inches(1), Inches(1)
    )
    p_circle.fill.solid()
    p_circle.fill.fore_color.rgb = RGBColor(52, 152, 219)
    p_circle.line.fill.background()
    
    p_text = shapes.add_textbox(Inches(7), Inches(3.7), Inches(1), Inches(0.6))
    tf = p_text.text_frame
    p_tf = tf.paragraphs[0]
    p_tf.text = "P"
    p_tf.font.name = "Arial"
    p_tf.font.size = Pt(36)
    p_tf.font.bold = True
    p_tf.font.color.rgb = RGBColor(255, 255, 255)
    p_tf.alignment = PP_ALIGN.CENTER
    
    # 连接线
    # C-A
    line1 = shapes.add_connector(
        1, Inches(4.7), Inches(2.3), Inches(2.8), Inches(3.7)
    )
    line1.line.color.rgb = RGBColor(189, 195, 199)
    line1.line.width = Pt(2)
    
    # C-P
    line2 = shapes.add_connector(
        1, Inches(5.3), Inches(2.3), Inches(7.2), Inches(3.7)
    )
    line2.line.color.rgb = RGBColor(189, 195, 199)
    line2.line.width = Pt(2)
    
    # A-P
    line3 = shapes.add_connector(
        1, Inches(3), Inches(4), Inches(7), Inches(4)
    )
    line3.line.color.rgb = RGBColor(189, 195, 199)
    line3.line.width = Pt(2)
    
    # 解释文字
    # Consistency
    c_desc = shapes.add_textbox(Inches(3.5), Inches(0.8), Inches(3), Inches(0.5))
    tf = c_desc.text_frame
    p = tf.paragraphs[0]
    p.text = u"一致性 (Consistency)"
    p.font.name = u"微软雅黑"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(231, 76, 60)
    p.alignment = PP_ALIGN.CENTER
    
    # Availability
    a_desc = shapes.add_textbox(Inches(0.5), Inches(3), Inches(2.5), Inches(0.5))
    tf = a_desc.text_frame
    p = tf.paragraphs[0]
    p.text = u"可用性 (Availability)"
    p.font.name = u"微软雅黑"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(46, 204, 113)
    p.alignment = PP_ALIGN.CENTER
    
    # Partition tolerance
    p_desc = shapes.add_textbox(Inches(6.5), Inches(3), Inches(3), Inches(0.5))
    tf = p_desc.text_frame
    p = tf.paragraphs[0]
    p.text = u"分区容错性 (Partition tolerance)"
    p.font.name = u"微软雅黑"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(52, 152, 219)
    p.alignment = PP_ALIGN.CENTER
    
    # 系统定位说明
    position_box = shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(3.5), Inches(2.5),
        Inches(3), Inches(0.8)
    )
    position_box.fill.solid()
    position_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
    position_box.line.color.rgb = RGBColor(41, 128, 185)
    position_box.line.width = Pt(2)
    
    position_text = shapes.add_textbox(Inches(3.5), Inches(2.6), Inches(3), Inches(0.6))
    tf = position_text.text_frame
    p = tf.paragraphs[0]
    p.text = u"本系统定位：CP系统"
    p.font.name = u"微软雅黑"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(41, 128, 185)
    p.alignment = PP_ALIGN.CENTER

def create_raft_architecture_slide(prs):
    """创建Raft算法架构图幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shapes = slide.shapes
    
    # 标题
    title_bg = shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(1)
    )
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = RGBColor(41, 128, 185)
    title_bg.line.fill.background()
    
    title_box = shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = u"Raft算法核心架构"
    p.font.name = u"微软雅黑"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Leader节点
    leader_box = shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(4), Inches(1.5),
        Inches(2), Inches(0.8)
    )
    leader_box.fill.solid()
    leader_box.fill.fore_color.rgb = RGBColor(231, 76, 60)
    leader_box.line.fill.background()
    
    leader_text = shapes.add_textbox(Inches(4), Inches(1.65), Inches(2), Inches(0.5))
    tf = leader_text.text_frame
    p = tf.paragraphs[0]
    p.text = u"Leader"
    p.font.name = u"微软雅黑"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Follower节点
    follower_positions = [(1.5, 3), (4, 3.5), (6.5, 3)]
    for i, (x, y) in enumerate(follower_positions):
        follower_box = shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y),
            Inches(2), Inches(0.8)
        )
        follower_box.fill.solid()
        follower_box.fill.fore_color.rgb = RGBColor(52, 152, 219)
        follower_box.line.fill.background()
        
        follower_text = shapes.add_textbox(Inches(x), Inches(y + 0.15), Inches(2), Inches(0.5))
        tf = follower_text.text_frame
        p = tf.paragraphs[0]
        p.text = u"Follower {}".format(i+1)
        p.font.name = u"微软雅黑"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
    
    # 添加箭头表示日志复制
    for x, y in follower_positions:
        arrow = shapes.add_connector(
            1, Inches(5), Inches(2.3), Inches(x + 1), Inches(y)
        )
        arrow.line.color.rgb = RGBColor(46, 204, 113)
        arrow.line.width = Pt(2)
    
    # 添加说明文字
    desc_box = shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(0.8))
    tf = desc_box.text_frame
    p = tf.paragraphs[0]
    p.text = u"Leader负责处理所有客户端请求，并将日志复制到Follower节点"
    p.font.name = u"微软雅黑"
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(52, 73, 94)
    p.alignment = PP_ALIGN.CENTER

def create_system_architecture_slide(prs):
    """创建系统架构图幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shapes = slide.shapes
    
    # 标题
    title_bg = shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(1)
    )
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = RGBColor(41, 128, 185)
    title_bg.line.fill.background()
    
    title_box = shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = u"系统架构设计"
    p.font.name = u"微软雅黑"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # 四层架构
    layers = [
        (u"接口层", u"API/SDK", RGBColor(52, 152, 219)),
        (u"服务层", u"业务逻辑", RGBColor(46, 204, 113)),
        (u"共识层", u"Raft实现", RGBColor(231, 76, 60)),
        (u"存储层", u"持久化", RGBColor(155, 89, 182))
    ]
    
    for i, (name, desc, color) in enumerate(layers):
        # 层级框
        layer_box = shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(2), Inches(1.5 + i * 0.9),
            Inches(6), Inches(0.8)
        )
        layer_box.fill.solid()
        layer_box.fill.fore_color.rgb = color
        layer_box.line.fill.background()
        
        # 层级名称
        name_text = shapes.add_textbox(Inches(2.2), Inches(1.6 + i * 0.9), Inches(2), Inches(0.6))
        tf = name_text.text_frame
        p = tf.paragraphs[0]
        p.text = name
        p.font.name = u"微软雅黑"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        
        # 层级描述
        desc_text = shapes.add_textbox(Inches(5.5), Inches(1.6 + i * 0.9), Inches(2.3), Inches(0.6))
        tf = desc_text.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.name = u"微软雅黑"
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.RIGHT
        
        # 连接箭头
        if i < len(layers) - 1:
            arrow = shapes.add_connector(
                1, Inches(5), Inches(2.3 + i * 0.9), Inches(5), Inches(2.3 + (i + 1) * 0.9)
            )
            arrow.line.color.rgb = RGBColor(127, 140, 141)
            arrow.line.width = Pt(2)

def create_performance_comparison_slide(prs):
    """创建性能对比图表幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shapes = slide.shapes
    
    # 标题
    title_bg = shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(1)
    )
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = RGBColor(41, 128, 185)
    title_bg.line.fill.background()
    
    title_box = shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = u"性能测试结果"
    p.font.name = u"微软雅黑"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # 创建柱状图
    chart_data = ChartData()
    chart_data.categories = ['48 keys', '512 keys', '4096 keys']
    chart_data.add_series('DDBR (秒)', (35.39, 39.32, 36.80))
    chart_data.add_series('Redis (秒)', (12.16, 12.47, 11.81))
    
    x, y, cx, cy = Inches(1), Inches(1.5), Inches(5), Inches(3)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    ).chart
    
    chart.has_title = True
    chart.chart_title.text_frame.text = u"10万次读操作耗时对比"
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(18)
    chart.chart_title.text_frame.paragraphs[0].font.bold = True
    
    # 添加分析文字
    analysis_box = shapes.add_textbox(Inches(6.5), Inches(1.5), Inches(3), Inches(3))
    tf = analysis_box.text_frame
    p = tf.paragraphs[0]
    p.text = u"性能分析"
    p.font.name = u"微软雅黑"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(41, 128, 185)
    
    analysis_points = [
        u"DDBR性能约为Redis的1/3",
        u"性能差距源于一致性保证",
        u"获得了强一致性和高可用性",
        u"符合CAP理论的权衡"
    ]
    
    for point in analysis_points:
        p = tf.add_paragraph()
        p.text = u"▸ " + point
        p.font.name = u"微软雅黑"
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(52, 73, 94)
        p.space_before = Pt(8)

def create_test_results_slide(prs):
    """创建测试结果展示幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shapes = slide.shapes
    
    # 标题
    title_bg = shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(1)
    )
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = RGBColor(41, 128, 185)
    title_bg.line.fill.background()
    
    title_box = shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = u"功能测试结果"
    p.font.name = u"微软雅黑"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # 测试项目卡片
    test_items = [
        (u"基本KV操作", u"✓ SET/GET/DELETE\n✓ 数据一致性验证\n✓ 并发操作测试", RGBColor(46, 204, 113)),
        (u"Leader选举", u"✓ 自动故障检测\n✓ 快速选举新Leader\n✓ 平均耗时<300ms", RGBColor(52, 152, 219)),
        (u"日志同步", u"✓ 实时数据同步\n✓ 冲突自动解决\n✓ 最终一致性保证", RGBColor(155, 89, 182)),
        (u"故障恢复", u"✓ 节点崩溃恢复\n✓ 网络分区处理\n✓ 数据完整性保证", RGBColor(231, 76, 60))
    ]
    
    for i, (title, content, color) in enumerate(test_items):
        # 创建测试卡片
        card = shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5 + (i % 2) * 4.8), Inches(1.5 + (i // 2) * 2),
            Inches(4.3), Inches(1.6)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(255, 255, 255)
        card.line.color.rgb = color
        card.line.width = Pt(3)
        
        # 标题
        title_text = shapes.add_textbox(
            Inches(0.7 + (i % 2) * 4.8), Inches(1.6 + (i // 2) * 2),
            Inches(3.9), Inches(0.4)
        )
        tf = title_text.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = u"微软雅黑"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = color
        
        # 内容
        content_text = shapes.add_textbox(
            Inches(0.7 + (i % 2) * 4.8), Inches(2 + (i // 2) * 2),
            Inches(3.9), Inches(1)
        )
        tf = content_text.text_frame
        p = tf.paragraphs[0]
        p.text = content
        p.font.name = u"微软雅黑"
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(52, 73, 94)

def create_achievements_slide(prs):
    """创建成果总结幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shapes = slide.shapes
    
    # 渐变背景
    add_gradient_background(slide, RGBColor(41, 128, 185), RGBColor(44, 62, 80))
    
    # 标题
    title_box = shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(0.8))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = u"研究成果总结"
    p.font.name = u"微软雅黑"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # 成果列表
    achievements = [
        u"✓ 深入研究了Raft共识算法的理论基础和实现原理",
        u"✓ 设计并实现了完整的分布式KV数据库系统",
        u"✓ 实现了Leader选举、日志复制、故障恢复等核心功能",
        u"✓ 完成了全面的功能测试和性能评估",
        u"✓ 在一致性和可用性之间找到了合理的平衡点",
        u"✓ 为分布式系统开发提供了实践参考"
    ]
    
    achievement_box = shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1), Inches(1.5),
        Inches(8), Inches(3)
    )
    achievement_box.fill.solid()
    achievement_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
    achievement_box.fill.transparency = 0.1
    achievement_box.line.fill.background()
    
    text_box = shapes.add_textbox(Inches(1.5), Inches(1.7), Inches(7), Inches(2.6))
    tf = text_box.text_frame
    
    for i, achievement in enumerate(achievements):
        if i > 0:
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[0]
        p.text = achievement
        p.font.name = u"微软雅黑"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.space_before = Pt(10)

def create_future_work_slide(prs):
    """创建未来展望幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shapes = slide.shapes
    
    # 标题
    title_bg = shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(1)
    )
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = RGBColor(41, 128, 185)
    title_bg.line.fill.background()
    
    title_box = shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = u"未来工作展望"
    p.font.name = u"微软雅黑"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # 未来工作方向
    future_items = [
        (u"性能优化", [u"ReadIndex机制", u"批量处理优化", u"并行复制"], RGBColor(52, 152, 219)),
        (u"功能扩展", [u"事务支持", u"数据分片", u"多数据类型"], RGBColor(46, 204, 113)),
        (u"可靠性提升", [u"成员变更优化", u"安全机制增强", u"跨数据中心"], RGBColor(231, 76, 60)),
        (u"可观测性", [u"监控系统", u"日志分析", u"可视化工具"], RGBColor(155, 89, 182))
    ]
    
    for i, (category, items, color) in enumerate(future_items):
        # 分类标题
        cat_box = shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5 + i * 2.4), Inches(1.5),
            Inches(2.2), Inches(0.6)
        )
        cat_box.fill.solid()
        cat_box.fill.fore_color.rgb = color
        cat_box.line.fill.background()
        
        cat_text = shapes.add_textbox(
            Inches(0.5 + i * 2.4), Inches(1.6),
            Inches(2.2), Inches(0.4)
        )
        tf = cat_text.text_frame
        p = tf.paragraphs[0]
        p.text = category
        p.font.name = u"微软雅黑"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        # 具体项目
        for j, item in enumerate(items):
            item_box = shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.6 + i * 2.4), Inches(2.3 + j * 0.6),
                Inches(2), Inches(0.5)
            )
            item_box.fill.solid()
            item_box.fill.fore_color.rgb = RGBColor(247, 249, 251)
            item_box.line.color.rgb = color
            item_box.line.width = Pt(1)
            
            item_text = shapes.add_textbox(
                Inches(0.6 + i * 2.4), Inches(2.35 + j * 0.6),
                Inches(2), Inches(0.4)
            )
            tf = item_text.text_frame
            p = tf.paragraphs[0]
            p.text = item
            p.font.name = u"微软雅黑"
            p.font.size = Pt(12)
            p.font.color.rgb = RGBColor(52, 73, 94)
            p.alignment = PP_ALIGN.CENTER

def create_thanks_slide(prs):
    """创建致谢幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 渐变背景
    add_gradient_background(slide, RGBColor(44, 62, 80), RGBColor(25, 52, 65))
    
    shapes = slide.shapes
    
    # 装饰圆圈
    for i in range(5):
        circle = shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(1 + i * 1.8), Inches(0.5),
            Inches(0.6), Inches(0.6)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = RGBColor(52, 152, 219)
        circle.fill.transparency = 0.5 + i * 0.1
        circle.line.fill.background()
    
    # 主标题
    title_box = shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1.5))
    text_frame = title_box.text_frame
    p = text_frame.paragraphs[0]
    p.text = u"感谢聆听"
    p.font.name = u"微软雅黑"
    p.font.size = Pt(64)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # 英文
    p = text_frame.add_paragraph()
    p.text = u"Thank You"
    p.font.name = "Arial"
    p.font.size = Pt(36)
    p.font.italic = True
    p.font.color.rgb = RGBColor(189, 195, 199)
    p.alignment = PP_ALIGN.CENTER
    
    # 分隔线
    line = shapes.add_connector(
        1, Inches(2), Inches(3.5), Inches(8), Inches(3.5)
    )
    line.line.color.rgb = RGBColor(255, 255, 255)
    line.line.width = Pt(2)
    
    # 底部信息
    info_box = shapes.add_textbox(Inches(1.5), Inches(4), Inches(7), Inches(1.2))
    text_frame = info_box.text_frame
    p = text_frame.paragraphs[0]
    p.text = u"请各位老师批评指正"
    p.font.name = u"微软雅黑"
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(236, 240, 241)
    p.alignment = PP_ALIGN.CENTER
    
    p = text_frame.add_paragraph()
    p.text = u"答辩人：张皓然  |  学号：202131061326"
    p.font.name = u"微软雅黑"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(189, 195, 199)
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(12)
    
    p = text_frame.add_paragraph()
    p.text = u"指导教师：蒋欣岑 讲师"
    p.font.name = u"微软雅黑"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(189, 195, 199)
    p.alignment = PP_ALIGN.CENTER

def create_enhanced_presentation():
    """创建增强版图文并茂的演示文稿"""
    prs = Presentation()
    
    # 设置16:9比例
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # 创建所有幻灯片
    create_title_slide(prs)
    create_outline_slide(prs)
    create_background_slide(prs)
    create_cap_slide(prs)
    create_raft_architecture_slide(prs)
    create_system_architecture_slide(prs)
    create_test_results_slide(prs)
    create_performance_comparison_slide(prs)
    create_achievements_slide(prs)
    create_future_work_slide(prs)
    create_thanks_slide(prs)
    
    return prs

def main():
    """主函数"""
    print(u"正在生成增强版图文并茂的毕业答辩PPT...")
    
    try:
        prs = create_enhanced_presentation()
        filename = u"毕业答辩_基于Raft算法的分布式数据库_增强版.pptx"
        prs.save(filename)
        print(u"\n✅ PPT生成成功！")
        print(u"📄 文件名：" + filename)
        print(u"📍 位置：" + os.path.abspath(filename))
        print(u"\n✨ 增强版特色：")
        print(u"  - 结合论文内容的详细技术展示")
        print(u"  - 丰富的图表和可视化元素")
        print(u"  - CAP理论三角形图解")
        print(u"  - Raft算法架构图")
        print(u"  - 系统四层架构图")
        print(u"  - 性能对比柱状图")
        print(u"  - 测试结果卡片式展示")
        print(u"  - 动态目录和装饰元素")
        print(u"  - 专业的配色和排版")
    except Exception as e:
        print(u"❌ 生成失败：" + str(e))
        import traceback
        traceback.print_exc()
        print(u"\n请确保：")
        print(u"1. 已安装python-pptx库：pip install python-pptx")
        print(u"2. Python版本支持中文处理")

if __name__ == "__main__":
    main() 