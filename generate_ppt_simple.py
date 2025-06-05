#!/usr/bin/env python
# -*- coding: utf-8 -*-

"""
简化版PPT生成脚本 - 兼容更多Python版本
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def create_presentation():
    """创建演示文稿"""
    prs = Presentation()
    
    # 1. 标题页
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = u"基于Raft算法的分布式数据库构建"
    subtitle.text = u"毕业设计答辩\n\n答辩人：张皓然\n学号：202131061326\n指导教师：蒋欣岑 讲师"
    
    # 2. 目录
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = u"目录"
    content.text = u"""1. 研究背景与目标
2. 相关理论基础
3. Raft算法核心机制
4. 系统设计与架构
5. 功能实现与测试
6. 性能分析
7. 总结与展望"""
    
    # 3. 研究背景
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = u"1. 研究背景"
    content.text = u"""大数据时代的挑战

数据爆炸式增长
- 传统单机数据库面临容量、性能瓶颈
- 分布式架构成为必然选择

核心技术难题
- 如何在分布式环境中保证数据一致性？
- 如何平衡一致性和可用性？"""
    
    # 4. 研究目标
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = u"2. 研究目标"
    content.text = u"""构建基于Raft的分布式KV数据库

主要目标：
- 实现强一致性的数据存储
- 保证系统高可用性
- 支持节点故障自动恢复
- 提供简洁的API接口

技术路线：
- 采用Raft共识算法
- Go语言开发
- Kitex RPC框架"""
    
    # 5. CAP理论
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = u"3. 理论基础 - CAP定理"
    content.text = u"""CAP三要素
- C (Consistency): 一致性
- A (Availability): 可用性
- P (Partition tolerance): 分区容错性

本系统定位：CP系统，优先保证数据一致性"""
    
    # 6. 为什么选择Raft
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = u"4. 为什么选择Raft？"
    content.text = u"""主流共识算法对比

Paxos: 复杂度高，难以理解和实现
Raft: 复杂度中等，简单易懂，容易实现
ZAB: 复杂度中等，理解和实现难度一般

Raft的优势：
- 算法清晰易懂
- 模块化设计
- 广泛的工业应用（etcd、TiKV等）"""
    
    # 7. Raft核心机制
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = u"5. Raft算法核心机制"
    content.text = u"""三种节点角色

Leader（领导者）
- 处理所有客户端请求
- 管理日志复制
- 发送心跳维持权威

Follower（跟随者）
- 被动接收日志
- 响应Leader请求

Candidate（候选者）
- 选举过程中的临时状态"""
    
    # 8. Leader选举
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = u"6. Leader选举机制"
    content.text = u"""选举触发条件
1. 系统初始化启动
2. Leader节点故障
3. 网络分区恢复

关键设计
- 随机超时：避免选票分散
- 任期机制：识别过期信息"""
    
    # 9. 日志复制
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = u"7. 日志复制流程"
    content.text = u"""复制步骤
1. 客户端发送写请求到Leader
2. Leader追加日志并复制到Follower
3. 多数节点确认后提交
4. 应用到状态机并返回结果

基于多数派的机制保证数据不丢失"""
    
    # 10. 系统架构
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = u"8. 系统架构设计"
    content.text = u"""分层架构
- 接口层 (API/SDK)
- 服务层 (业务逻辑)
- 共识层 (Raft实现)
- 存储层 (持久化)

技术栈
- 开发语言: Go
- RPC框架: Kitex + Thrift
- 存储引擎: 内存KV + WAL"""
    
    # 11. 功能实现
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = u"9. 核心功能实现"
    content.text = u"""基础KV操作
- SET(key, value) - 设置键值对
- GET(key) -> value - 获取值
- DELETE(key) - 删除键

集群管理功能
- 节点动态加入/退出
- 自动故障检测与恢复
- 日志同步与一致性保证"""
    
    # 12. Leader选举测试
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = u"10. 功能测试展示"
    content.text = u"""1. Leader选举测试

测试场景：关闭Leader节点
测试结果：Follower 2成功当选新Leader

故障恢复时间：秒级完成"""
    
    # 13. 日志同步测试
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = u"11. 功能测试展示"
    content.text = u"""2. 日志同步测试

测试步骤：
1. Leader执行SET操作
2. 观察Follower日志输出
3. 验证数据一致性

结果：所有节点数据保持一致"""
    
    # 14. 性能分析
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = u"12. 性能测试分析"
    content.text = u"""与Redis性能对比（10万次读操作）

负载规模 | DDBR耗时 | Redis耗时 | 性能比
48 keys  | 35.39s   | 12.16s    | 34%
512 keys | 39.32s   | 12.47s    | 32%
4096 keys| 36.80s   | 11.81s    | 32%

分析
- 性能差距主要源于一致性保证开销
- 但获得了强一致性和高可用性
- 符合CAP理论的权衡"""
    
    # 15. 系统优势
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = u"13. 系统优势与创新点"
    content.text = u"""技术优势
- 强一致性保证：基于Raft算法
- 高可用性：自动故障恢复
- 模块化设计：易于扩展维护
- 高性能通信：Kitex RPC框架

创新点
- 实现了完整的Raft核心机制
- 优化了日志复制性能
- 提供了多种一致性级别的读取策略"""
    
    # 16. 总结
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = u"14. 总结"
    content.text = u"""完成的工作
- 深入研究了Raft共识算法原理
- 设计并实现了分布式KV数据库系统
- 实现了Leader选举、日志复制等核心功能
- 完成了功能测试和性能评估

达成的目标
- 构建了一个可用的分布式存储系统
- 在一致性和可用性间找到平衡
- 为分布式系统开发提供了实践参考"""
    
    # 17. 展望
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = u"15. 未来展望"
    content.text = u"""性能优化
- 实现ReadIndex和LeaderLease机制
- 引入批量处理和并行复制
- 优化网络传输和序列化

功能扩展
- 增加事务支持
- 实现数据分片
- 支持更多数据类型

可靠性提升
- 完善成员变更机制
- 增强安全机制
- 支持跨数据中心部署"""
    
    # 18. 致谢
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = u"谢谢！"
    content.text = u"""请各位老师批评指正

答辩人：张皓然
学号：202131061326
指导教师：蒋欣岑 讲师"""
    
    return prs

def main():
    """主函数"""
    print("正在生成毕业答辩PPT (简化版)...")
    
    try:
        prs = create_presentation()
        filename = "毕业答辩_基于Raft算法的分布式数据库构建_简化版.pptx"
        prs.save(filename)
        print("PPT生成成功！")
        print("文件名：" + filename)
        print("位置：" + os.path.abspath(filename))
    except Exception as e:
        print("生成失败：" + str(e))
        print("请确保已安装python-pptx库")
        print("安装命令：pip install python-pptx")

if __name__ == "__main__":
    main() 