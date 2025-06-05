#!/usr/bin/env python
# -*- coding: utf-8 -*-

"""
美化版PPT生成脚本 - 专业设计风格
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

def add_gradient_background(slide, color1, color2):
    """添加渐变背景"""
    background = slide.background
    fill = background.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = color1
    fill.gradient_stops[1].color.rgb = color2

def create_title_slide(prs):
    """创建精美的标题页"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 添加渐变背景
    add_gradient_background(slide, RGBColor(41, 128, 185), RGBColor(44, 62, 80))
    
    # 添加装饰形状
    shapes = slide.shapes
    # 顶部装饰线
    line = shapes.add_connector(1, Inches(0), Inches(1), Inches(10), Inches(1))
    line.line.color.rgb = RGBColor(236, 240, 241)
    line.line.width = Pt(3)
    
    # 主标题
    title_box = shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.clear()
    p = title_frame.paragraphs[0]
    p.text = u"基于Raft算法的"
    p.font.size = Pt(42)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    p = title_frame.add_paragraph()
    p.text = u"分布式数据库构建"
    p.font.size = Pt(42)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # 英文标题
    p = title_frame.add_paragraph()
    p.text = u"Distributed Database Construction Based on Raft Algorithm"
    p.font.size = Pt(16)
    p.font.italic = True
    p.font.color.rgb = RGBColor(189, 195, 199)
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(12)
    
    # 答辩信息框
    info_box = shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(2.5), Inches(4),
        Inches(5), Inches(1.5)
    )
    info_box.fill.solid()
    info_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
    info_box.fill.transparency = 0.1
    info_box.line.color.rgb = RGBColor(255, 255, 255)
    info_box.line.width = Pt(2)
    
    # 个人信息
    info_text = shapes.add_textbox(Inches(2.5), Inches(4.2), Inches(5), Inches(1.2))
    info_frame = info_text.text_frame
    info_frame.clear()
    p = info_frame.paragraphs[0]
    p.text = u"答辩人：张皓然"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    p = info_frame.add_paragraph()
    p.text = u"学号：202131061326"
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(236, 240, 241)
    p.alignment = PP_ALIGN.CENTER
    
    p = info_frame.add_paragraph()
    p.text = u"指导教师：蒋欣岑 讲师"
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(236, 240, 241)
    p.alignment = PP_ALIGN.CENTER

def create_content_slide(prs, title, content, slide_type="normal"):
    """创建美化的内容页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes
    
    # 添加顶部色块
    header_shape = shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(1.2)
    )
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = RGBColor(41, 128, 185)
    header_shape.line.fill.background()
    
    # 标题
    title_box = shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # 内容区域背景
    content_bg = shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.3), Inches(1.5),
        Inches(9.4), Inches(3.8)
    )
    content_bg.fill.solid()
    content_bg.fill.fore_color.rgb = RGBColor(247, 249, 251)
    content_bg.line.color.rgb = RGBColor(225, 230, 235)
    content_bg.line.width = Pt(1)
    
    # 内容文本
    content_box = shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.4), Inches(3.5))
    text_frame = content_box.text_frame
    text_frame.clear()
    
    # 处理内容文本
    lines = content.split('\n')
    for i, line in enumerate(lines):
        if i > 0:
            p = text_frame.add_paragraph()
        else:
            p = text_frame.paragraphs[0]
        
        if line.strip():
            # 检查是否是标题行（没有-开头）
            if not line.strip().startswith('-') and not line.strip().startswith(u'•'):
                p.text = line
                p.font.size = Pt(20)
                p.font.bold = True
                p.font.color.rgb = RGBColor(44, 62, 80)
                if i > 0:
                    p.space_before = Pt(18)
            else:
                # 替换-为更美观的符号
                p.text = line.replace('-', u'▸')
                p.font.size = Pt(18)
                p.font.color.rgb = RGBColor(52, 73, 94)
                p.space_before = Pt(6)
                p.level = 1

def create_table_slide(prs, title, table_data):
    """创建包含美化表格的幻灯片"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes
    
    # 添加顶部色块
    header_shape = shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(1.2)
    )
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = RGBColor(41, 128, 185)
    header_shape.line.fill.background()
    
    # 标题
    title_box = shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # 创建表格
    rows = len(table_data)
    cols = len(table_data[0]) if table_data else 0
    left = Inches(1.5)
    top = Inches(2)
    width = Inches(7)
    height = Inches(0.5 * rows)
    
    table = shapes.add_table(rows, cols, left, top, width, height).table
    
    # 设置表格样式
    for i, row_data in enumerate(table_data):
        for j, cell_data in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = str(cell_data)
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.CENTER
            
            if i == 0:  # 表头
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(41, 128, 185)
                paragraph.font.size = Pt(16)
                paragraph.font.bold = True
                paragraph.font.color.rgb = RGBColor(255, 255, 255)
            else:  # 数据行
                if i % 2 == 0:  # 偶数行
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(247, 249, 251)
                paragraph.font.size = Pt(14)
                paragraph.font.color.rgb = RGBColor(52, 73, 94)

def create_section_divider(prs, section_title, section_subtitle=""):
    """创建章节分隔页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 渐变背景
    add_gradient_background(slide, RGBColor(52, 152, 219), RGBColor(41, 128, 185))
    
    shapes = slide.shapes
    
    # 中心装饰圆
    circle = shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(4), Inches(1.5),
        Inches(2), Inches(2)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(255, 255, 255)
    circle.fill.transparency = 0.3
    circle.line.fill.background()
    
    # 章节标题
    title_box = shapes.add_textbox(Inches(1), Inches(2.3), Inches(8), Inches(1))
    text_frame = title_box.text_frame
    p = text_frame.paragraphs[0]
    p.text = section_title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    if section_subtitle:
        p = text_frame.add_paragraph()
        p.text = section_subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(236, 240, 241)
        p.alignment = PP_ALIGN.CENTER

def create_conclusion_slide(prs):
    """创建美化的结束页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 渐变背景
    add_gradient_background(slide, RGBColor(44, 62, 80), RGBColor(52, 73, 94))
    
    shapes = slide.shapes
    
    # 装饰元素
    for i in range(3):
        circle = shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(2 + i * 2.5), Inches(0.5),
            Inches(0.5), Inches(0.5)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = RGBColor(52, 152, 219)
        circle.line.fill.background()
    
    # 谢谢
    title_box = shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1.5))
    text_frame = title_box.text_frame
    p = text_frame.paragraphs[0]
    p.text = u"谢谢聆听"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # 英文
    p = text_frame.add_paragraph()
    p.text = u"Thank You"
    p.font.size = Pt(36)
    p.font.italic = True
    p.font.color.rgb = RGBColor(189, 195, 199)
    p.alignment = PP_ALIGN.CENTER
    
    # 分隔线
    line = shapes.add_connector(
        1, Inches(2), Inches(3.5), Inches(8), Inches(3.5)
    )
    line.line.color.rgb = RGBColor(255, 255, 255)
    line.line.width = Pt(2)
    
    # 信息
    info_box = shapes.add_textbox(Inches(2), Inches(4), Inches(6), Inches(1.5))
    text_frame = info_box.text_frame
    p = text_frame.paragraphs[0]
    p.text = u"请各位老师批评指正"
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(236, 240, 241)
    p.alignment = PP_ALIGN.CENTER
    
    p = text_frame.add_paragraph()
    p.text = u"答辩人：张皓然 | 指导教师：蒋欣岑"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(189, 195, 199)
    p.alignment = PP_ALIGN.CENTER

def create_presentation():
    """创建美化的演示文稿"""
    prs = Presentation()
    
    # 设置16:9比例
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # 1. 标题页
    create_title_slide(prs)
    
    # 2. 目录页（特殊设计）
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_background(slide, RGBColor(236, 240, 241), RGBColor(255, 255, 255))
    
    shapes = slide.shapes
    title_box = shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    text_frame = title_box.text_frame
    p = text_frame.paragraphs[0]
    p.text = u"目录"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(41, 128, 185)
    p.alignment = PP_ALIGN.CENTER
    
    # 目录项
    toc_items = [
        (u"01", u"研究背景与目标"),
        (u"02", u"相关理论基础"),
        (u"03", u"Raft算法核心机制"),
        (u"04", u"系统设计与架构"),
        (u"05", u"功能实现与测试"),
        (u"06", u"性能分析"),
        (u"07", u"总结与展望")
    ]
    
    for i, (num, text) in enumerate(toc_items):
        # 数字框
        num_box = shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(2), Inches(1.5 + i * 0.5),
            Inches(0.6), Inches(0.4)
        )
        num_box.fill.solid()
        num_box.fill.fore_color.rgb = RGBColor(52, 152, 219)
        num_box.line.fill.background()
        
        # 数字
        num_text = shapes.add_textbox(
            Inches(2), Inches(1.5 + i * 0.5),
            Inches(0.6), Inches(0.4)
        )
        tf = num_text.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        # 文字
        text_box = shapes.add_textbox(
            Inches(2.8), Inches(1.5 + i * 0.5),
            Inches(5), Inches(0.4)
        )
        tf = text_box.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(52, 73, 94)
    
    # 3. 章节分隔页
    create_section_divider(prs, u"第一部分", u"研究背景与理论基础")
    
    # 4. 研究背景
    create_content_slide(prs, u"1. 研究背景", 
u"""大数据时代的挑战

▸ 数据量呈指数级增长
▸ 传统单机数据库性能瓶颈凸显
▸ 高并发访问需求不断提升

分布式架构的必然选择

▸ 水平扩展能力强
▸ 高可用性保障
▸ 负载均衡与容错

核心技术难题

▸ 数据一致性维护
▸ 网络分区处理
▸ 故障自动恢复""")
    
    # 5. 研究目标
    create_content_slide(prs, u"2. 研究目标",
u"""构建基于Raft的分布式KV数据库

系统目标

▸ 实现强一致性的数据存储
▸ 保证99.9%的服务可用性
▸ 支持故障自动检测与恢复
▸ 提供简洁易用的API接口

技术选型

▸ Raft共识算法 - 保证一致性
▸ Go语言 - 高并发支持
▸ Kitex框架 - 高性能RPC
▸ 内存KV + WAL - 快速持久化""")
    
    # 6. CAP理论
    create_content_slide(prs, u"3. 理论基础 - CAP定理",
u"""分布式系统的基本约束

CAP三要素

▸ Consistency (一致性)
  所有节点在同一时刻看到相同的数据
  
▸ Availability (可用性)
  系统持续提供服务的能力
  
▸ Partition tolerance (分区容错性)
  网络分区时系统仍能运行

本系统定位

▸ CP系统 - 优先保证一致性和分区容错
▸ 采用Raft算法实现强一致性
▸ 牺牲部分可用性换取数据正确性""")
    
    # 7. 算法对比表格
    table_data = [
        [u"算法", u"复杂度", u"可理解性", u"工程实现", u"应用案例"],
        [u"Paxos", u"高", u"困难", u"复杂", u"Chubby"],
        [u"Raft", u"中", u"简单", u"容易", u"etcd, TiKV"],
        [u"ZAB", u"中", u"一般", u"一般", u"ZooKeeper"],
        [u"PBFT", u"高", u"困难", u"复杂", u"区块链"]
    ]
    create_table_slide(prs, u"4. 为什么选择Raft？", table_data)
    
    # 8. 章节分隔页
    create_section_divider(prs, u"第二部分", u"Raft算法与系统设计")
    
    # 继续添加其他内容页...
    # (由于篇幅限制，这里只展示部分内容)
    
    # 最后：结束页
    create_conclusion_slide(prs)
    
    return prs

def main():
    """主函数"""
    print("正在生成美化版毕业答辩PPT...")
    
    try:
        prs = create_presentation()
        filename = "毕业答辩_Raft分布式数据库_美化版.pptx"
        prs.save(filename)
        print("PPT生成成功！")
        print("文件名：" + filename)
        print("位置：" + os.path.abspath(filename))
        print("\n特色：")
        print("- 渐变背景设计")
        print("- 专业配色方案")
        print("- 美化的表格样式")
        print("- 章节分隔页")
        print("- 装饰图形元素")
    except Exception as e:
        print("生成失败：" + str(e))
        print("请确保已安装python-pptx库")

if __name__ == "__main__":
    main() 