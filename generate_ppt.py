#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
生成毕业答辩PPT的Python脚本
需要安装：pip install python-pptx pillow
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

def create_title_slide(prs):
    """创建标题页"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 设置深蓝色背景
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(44, 62, 80)  # #2c3e50
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "基于Raft算法的分布式数据库构建"
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 副标题
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "毕业设计答辩"
    subtitle_frame.paragraphs[0].font.size = Pt(24)
    subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(236, 240, 241)
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 个人信息
    info_box = slide.shapes.add_textbox(Inches(3), Inches(5), Inches(4), Inches(2))
    info_frame = info_box.text_frame
    info_frame.text = "答辩人：张皓然\n学号：202131061326\n指导教师：蒋欣岑 讲师"
    for paragraph in info_frame.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        paragraph.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, content, has_table=False, table_data=None):
    """添加内容页"""
    slide_layout = prs.slide_layouts[5]  # 标题和内容
    slide = prs.slides.add_slide(slide_layout)
    
    # 设置标题
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(28)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)
    
    # 添加内容
    if not has_table:
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
        content_frame = content_box.text_frame
        content_frame.text = content
        
        for paragraph in content_frame.paragraphs:
            paragraph.font.size = Pt(18)
            paragraph.space_before = Pt(6)
            paragraph.space_after = Pt(6)
    else:
        # 添加表格
        if table_data:
            rows, cols = len(table_data), len(table_data[0])
            left = Inches(2)
            top = Inches(2)
            width = Inches(6)
            height = Inches(0.8 * rows)
            
            table = slide.shapes.add_table(rows, cols, left, top, width, height).table
            
            # 填充表格数据
            for i, row in enumerate(table_data):
                for j, cell in enumerate(row):
                    table.cell(i, j).text = str(cell)
                    table.cell(i, j).text_frame.paragraphs[0].font.size = Pt(14)
                    if i == 0:  # 表头
                        table.cell(i, j).fill.solid()
                        table.cell(i, j).fill.fore_color.rgb = RGBColor(52, 73, 94)
                        table.cell(i, j).text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                        table.cell(i, j).text_frame.paragraphs[0].font.bold = True

def create_presentation():
    """创建完整的演示文稿"""
    prs = Presentation()
    
    # 设置16:9比例
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # 1. 标题页
    create_title_slide(prs)
    
    # 2. 目录页
    add_content_slide(prs, "目录", 
"""1. 研究背景与目标
2. 相关理论基础
3. Raft算法核心机制
4. 系统设计与架构
5. 功能实现与测试
6. 性能分析
7. 总结与展望""")
    
    # 3. 研究背景
    add_content_slide(prs, "1. 研究背景",
"""大数据时代的挑战

📈 数据爆炸式增长
• 传统单机数据库面临容量、性能瓶颈
• 分布式架构成为必然选择

🔧 核心技术难题
• 如何在分布式环境中保证数据一致性？
• 如何平衡一致性和可用性？""")
    
    # 4. 研究目标
    add_content_slide(prs, "2. 研究目标",
"""构建基于Raft的分布式KV数据库

主要目标：
✅ 实现强一致性的数据存储
✅ 保证系统高可用性
✅ 支持节点故障自动恢复
✅ 提供简洁的API接口

技术路线：
• 采用Raft共识算法
• Go语言开发
• Kitex RPC框架""")
    
    # 5. CAP理论
    add_content_slide(prs, "3. 理论基础 - CAP定理",
"""CAP三要素
• C (Consistency): 一致性
• A (Availability): 可用性
• P (Partition tolerance): 分区容错性

本系统定位：CP系统，优先保证数据一致性""")
    
    # 6. 为什么选择Raft - 带表格
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "4. 为什么选择Raft？"
    
    # 添加表格
    table_data = [
        ["算法", "复杂度", "可理解性", "工程实现"],
        ["Paxos", "高", "困难", "复杂"],
        ["Raft", "中", "简单", "容易"],
        ["ZAB", "中", "一般", "一般"]
    ]
    
    # 先添加说明文字
    text_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1))
    text_frame = text_box.text_frame
    text_frame.text = "主流共识算法对比"
    text_frame.paragraphs[0].font.size = Pt(20)
    text_frame.paragraphs[0].font.bold = True
    
    # 添加表格
    rows, cols = len(table_data), len(table_data[0])
    left = Inches(2)
    top = Inches(2.5)
    width = Inches(6)
    height = Inches(2)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    for i, row in enumerate(table_data):
        for j, cell in enumerate(row):
            table.cell(i, j).text = str(cell)
            table.cell(i, j).text_frame.paragraphs[0].font.size = Pt(14)
            if i == 0:
                table.cell(i, j).fill.solid()
                table.cell(i, j).fill.fore_color.rgb = RGBColor(52, 73, 94)
                table.cell(i, j).text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                table.cell(i, j).text_frame.paragraphs[0].font.bold = True
    
    # 添加优势说明
    text_box2 = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(1))
    text_frame2 = text_box2.text_frame
    text_frame2.text = "Raft的优势：✨ 算法清晰易懂  ✨ 模块化设计  ✨ 广泛的工业应用"
    text_frame2.paragraphs[0].font.size = Pt(16)
    
    # 7. Raft核心机制
    add_content_slide(prs, "5. Raft算法核心机制",
"""三种节点角色

Leader（领导者）
  ├─ 处理所有客户端请求
  ├─ 管理日志复制
  └─ 发送心跳维持权威

Follower（跟随者）
  ├─ 被动接收日志
  └─ 响应Leader请求

Candidate（候选者）
  └─ 选举过程中的临时状态""")
    
    # 8. Leader选举
    add_content_slide(prs, "6. Leader选举机制",
"""选举触发条件
1. 系统初始化启动
2. Leader节点故障
3. 网络分区恢复

关键设计
• 随机超时：避免选票分散
• 任期机制：识别过期信息""")
    
    # 9. 日志复制
    add_content_slide(prs, "7. 日志复制流程",
"""复制步骤
1. 客户端发送写请求到Leader
2. Leader追加日志并复制到Follower
3. 多数节点确认后提交
4. 应用到状态机并返回结果

基于多数派的机制保证数据不丢失""")
    
    # 10. 系统架构
    add_content_slide(prs, "8. 系统架构设计",
"""分层架构

┌─────────────────────────────┐
│      接口层 (API/SDK)        │
├─────────────────────────────┤
│      服务层 (业务逻辑)       │
├─────────────────────────────┤
│      共识层 (Raft实现)       │
├─────────────────────────────┤
│      存储层 (持久化)         │
└─────────────────────────────┘

技术栈
• 开发语言: Go
• RPC框架: Kitex + Thrift
• 存储引擎: 内存KV + WAL""")
    
    # 11. 功能实现
    add_content_slide(prs, "9. 核心功能实现",
"""基础KV操作
• SET(key, value) - 设置键值对
• GET(key) -> value - 获取值
• DELETE(key) - 删除键

集群管理功能
✅ 节点动态加入/退出
✅ 自动故障检测与恢复
✅ 日志同步与一致性保证""")
    
    # 12. Leader选举测试
    add_content_slide(prs, "10. 功能测试展示",
"""1. Leader选举测试

测试场景：关闭Leader节点
测试结果：Follower 2成功当选新Leader

故障恢复时间：秒级完成""")
    
    # 13. 日志同步测试
    add_content_slide(prs, "11. 功能测试展示",
"""2. 日志同步测试

测试步骤：
1. Leader执行SET操作
2. 观察Follower日志输出
3. 验证数据一致性

结果：所有节点数据保持一致 ✅""")
    
    # 14. 性能分析 - 带表格
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "12. 性能测试分析"
    
    # 添加说明
    text_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(0.5))
    text_frame = text_box.text_frame
    text_frame.text = "与Redis性能对比（10万次读操作）"
    text_frame.paragraphs[0].font.size = Pt(18)
    text_frame.paragraphs[0].font.bold = True
    
    # 添加表格
    perf_data = [
        ["负载规模", "DDBR耗时", "Redis耗时", "性能比"],
        ["48 keys", "35.39s", "12.16s", "34%"],
        ["512 keys", "39.32s", "12.47s", "32%"],
        ["4096 keys", "36.80s", "11.81s", "32%"]
    ]
    
    rows, cols = len(perf_data), len(perf_data[0])
    left = Inches(2)
    top = Inches(2.2)
    width = Inches(6)
    height = Inches(1.8)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    for i, row in enumerate(perf_data):
        for j, cell in enumerate(row):
            table.cell(i, j).text = str(cell)
            table.cell(i, j).text_frame.paragraphs[0].font.size = Pt(14)
            if i == 0:
                table.cell(i, j).fill.solid()
                table.cell(i, j).fill.fore_color.rgb = RGBColor(52, 73, 94)
                table.cell(i, j).text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                table.cell(i, j).text_frame.paragraphs[0].font.bold = True
    
    # 添加分析
    text_box2 = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1.2))
    text_frame2 = text_box2.text_frame
    text_frame2.text = """分析
• 性能差距主要源于一致性保证开销
• 但获得了强一致性和高可用性
• 符合CAP理论的权衡"""
    for paragraph in text_frame2.paragraphs:
        paragraph.font.size = Pt(16)
    
    # 15. 系统优势
    add_content_slide(prs, "13. 系统优势与创新点",
"""技术优势
🚀 强一致性保证：基于Raft算法
🛡️ 高可用性：自动故障恢复
📦 模块化设计：易于扩展维护
⚡ 高性能通信：Kitex RPC框架

创新点
• 实现了完整的Raft核心机制
• 优化了日志复制性能
• 提供了多种一致性级别的读取策略""")
    
    # 16. 总结
    add_content_slide(prs, "14. 总结",
"""完成的工作
✅ 深入研究了Raft共识算法原理
✅ 设计并实现了分布式KV数据库系统
✅ 实现了Leader选举、日志复制等核心功能
✅ 完成了功能测试和性能评估

达成的目标
• 构建了一个可用的分布式存储系统
• 在一致性和可用性间找到平衡
• 为分布式系统开发提供了实践参考""")
    
    # 17. 展望
    add_content_slide(prs, "15. 未来展望",
"""性能优化
• 实现ReadIndex和LeaderLease机制
• 引入批量处理和并行复制
• 优化网络传输和序列化

功能扩展
• 增加事务支持
• 实现数据分片
• 支持更多数据类型

可靠性提升
• 完善成员变更机制
• 增强安全机制
• 支持跨数据中心部署""")
    
    # 18. 致谢页
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 设置深蓝色背景
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(44, 62, 80)
    
    # 谢谢
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "谢谢！"
    title_frame.paragraphs[0].font.size = Pt(48)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 副标题
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(9), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "请各位老师批评指正"
    subtitle_frame.paragraphs[0].font.size = Pt(24)
    subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(236, 240, 241)
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 个人信息
    info_box = slide.shapes.add_textbox(Inches(3), Inches(4.5), Inches(4), Inches(1))
    info_frame = info_box.text_frame
    info_frame.text = "答辩人：张皓然\n学号：202131061326\n指导教师：蒋欣岑 讲师"
    for paragraph in info_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        paragraph.alignment = PP_ALIGN.CENTER
    
    return prs

def main():
    """主函数"""
    print("正在生成毕业答辩PPT...")
    
    try:
        prs = create_presentation()
        filename = "毕业答辩_基于Raft算法的分布式数据库构建.pptx"
        prs.save(filename)
        print("✅ PPT生成成功！文件名：" + filename)
        print("📍 文件位置：" + os.path.abspath(filename))
    except Exception as e:
        print("❌ 生成失败：" + str(e))
        print("请确保已安装必要的库：pip install python-pptx")

if __name__ == "__main__":
    main() 